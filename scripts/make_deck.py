#!/usr/bin/env python3
"""Build a high-design deck from HTML, then export PDF and editable PPTX.

Outputs in deliverables/:
- final_presentation.html
- final_presentation.pdf
- final_presentation.pptx
"""

from __future__ import annotations

import asyncio
import base64
import csv
import html
import json
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
DELIVERABLES = ROOT / "deliverables"
REPORTS = ROOT / "reports"
FIGURES = REPORTS / "figures"

HTML_OUT = DELIVERABLES / "final_presentation.html"
PDF_OUT = DELIVERABLES / "final_presentation.pdf"
PPTX_OUT = DELIVERABLES / "final_presentation.pptx"

REPO_LINK = "https://github.com/othmane-zizi-pro/nwa"


def _read_csv_rows(path: Path) -> List[Dict[str, str]]:
    if not path.exists():
        return []
    with path.open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def load_dynamic_metrics() -> Dict[str, str]:
    metrics = {
        "best_model": "Lasso Regression (Tuned)",
        "rmse": "1756.13",
        "mae": "545.35",
        "r2": "0.810",
        "lift_10": "5.55",
        "max_psi": "0.035",
        "ate_lrs": "-0.1785",
        "ate_xgbt": "-0.1673",
    }

    model_rows = _read_csv_rows(REPORTS / "model_results.csv")
    if model_rows:
        top = model_rows[0]
        metrics["best_model"] = top.get("", top.get("Model", metrics["best_model"]))
        metrics["rmse"] = f"{float(top.get('RMSE', 0.0)):.2f}"
        metrics["mae"] = f"{float(top.get('MAE', 0.0)):.2f}"
        metrics["r2"] = f"{float(top.get('R2', 0.0)):.3f}"

    meta = REPORTS / "pipeline_metadata.json"
    if meta.exists():
        payload = json.loads(meta.read_text(encoding="utf-8"))
        lift_10 = payload.get("lift_stats", {}).get("Lift@10%")
        if lift_10 is not None:
            metrics["lift_10"] = f"{float(lift_10):.2f}"

    drift_rows = _read_csv_rows(REPORTS / "monitoring_input_drift.csv")
    if drift_rows:
        max_psi = max(float(r.get("PSI", 0.0)) for r in drift_rows)
        metrics["max_psi"] = f"{max_psi:.3f}"

    causal_rows = _read_csv_rows(REPORTS / "causal_ate_results.csv")
    by_model = {r.get("Model", ""): r for r in causal_rows}
    if "LRSRegressor" in by_model:
        metrics["ate_lrs"] = f"{float(by_model['LRSRegressor']['ATE_log1p']):.4f}"
    if "XGBTRegressor" in by_model:
        metrics["ate_xgbt"] = f"{float(by_model['XGBTRegressor']['ATE_log1p']):.4f}"

    return metrics


def img_data_uri(path: Path) -> str:
    if not path.exists():
        return ""
    ext = path.suffix.lower().replace(".", "") or "png"
    payload = base64.b64encode(path.read_bytes()).decode("utf-8")
    return f"data:image/{ext};base64,{payload}"


def build_slides(metrics: Dict[str, str]) -> List[Dict[str, object]]:
    return [
        {
            "layout": "hero",
            "kicker": "INSY 674 Final Project",
            "title": "Customer Lifetime Value Prediction",
            "subtitle": (
                "Predictive modeling, explainability, causal inference, and monitoring "
                "for an end-to-end enterprise data science workflow."
            ),
            "footer": (
                "Othmane Zizi (261255341) | Fares Joni (261254593) | "
                "Tanmay Giri (261272443) | " + REPO_LINK
            ),
        },
        {
            "layout": "split",
            "kicker": "Context",
            "title": "Business Context and Objective (5.9.1)",
            "left_title": "Problem",
            "left_bullets": [
                "Marketing budgets are often distributed uniformly across unequal customers.",
                "High-value customers can churn without proactive intervention.",
                "Revenue concentration means ranking quality drives business impact.",
            ],
            "right_title": "Objective",
            "right_bullets": [
                "Predict 6-month CLV per customer from historical transaction behavior.",
                "Use model output to prioritize retention and campaign spend.",
                "Track value capture with lift and production monitoring.",
            ],
        },
        {
            "layout": "bullets",
            "kicker": "Hypothesis",
            "title": "Hypotheses and Testing Framing (5.9.2)",
            "bullets": [
                "H1: RFM + behavioral features can predict future CLV.",
                "H2: Regularized linear models generalize better than tree models on this feature set.",
                "H3: Causal uplift analysis identifies who should receive treatment.",
                "Null framing: treatment has no measurable effect on log1p(CLV).",
            ],
        },
        {
            "layout": "gallery",
            "kicker": "EDA",
            "title": "Exploration Visuals: Coverage and Quality (5.3)",
            "summary": [
                "Initial EDA covered time trends, missingness, geography, and transactional behavior.",
                "These checks informed cleaning decisions before feature engineering.",
            ],
            "images": [
                {"path": FIGURES / "daily_transactions.png", "caption": "Daily transaction volume"},
                {"path": FIGURES / "missing_values.png", "caption": "Missing values profile"},
                {"path": FIGURES / "top_countries.png", "caption": "Top purchasing countries"},
                {"path": FIGURES / "price_quantity_dist.png", "caption": "Price and quantity distribution"},
            ],
        },
        {
            "layout": "gallery",
            "kicker": "Features",
            "title": "Feature and Target Distributions (5.4)",
            "summary": [
                "RFM and behavioral features show heavy skew and strong concentration effects.",
                "Correlation and distribution checks shaped model family and scaling decisions.",
            ],
            "images": [
                {"path": FIGURES / "rfm_distributions.png", "caption": "RFM feature distributions"},
                {"path": FIGURES / "clv_distribution.png", "caption": "Future CLV distribution"},
                {"path": FIGURES / "customer_frequency.png", "caption": "Customer purchase frequency"},
                {"path": FIGURES / "correlation_matrix.png", "caption": "Feature correlation matrix"},
            ],
        },
        {
            "layout": "gallery",
            "kicker": "Segmentation",
            "title": "Segmentation and Behavioral Clusters",
            "summary": [
                "Segmentation complements prediction by converting scores into action buckets.",
                "Profiles and cluster diagnostics guide campaign strategy design.",
            ],
            "images": [
                {"path": FIGURES / "segment_distribution.png", "caption": "Segment size distribution"},
                {"path": FIGURES / "segment_profiles.png", "caption": "Segment profile comparison"},
                {"path": FIGURES / "rfm_segments.png", "caption": "RFM segment map"},
                {"path": FIGURES / "elbow_method.png", "caption": "K-means elbow method"},
            ],
        },
        {
            "layout": "bullets",
            "kicker": "Modeling",
            "title": "Modeling Approach and Evaluation (5.5, 5.6)",
            "bullets": [
                "Train/val/test split: 80/10/10 at customer level; validation for tuning, test as final holdout.",
                "Compared Linear, Ridge, Lasso, Random Forest, Gradient Boosting, XGBoost.",
                "Metrics: RMSE, MAE, R2, MAPE and lift for business prioritization.",
                "Randomized search used for fine-tuning top candidates.",
            ],
        },
        {
            "layout": "table",
            "kicker": "Results",
            "title": "Predictive Performance (5.9.5)",
            "headers": ["Model", "RMSE", "MAE", "R2"],
            "rows": [
                [metrics["best_model"], metrics["rmse"], metrics["mae"], metrics["r2"]],
                ["Lasso Regression", "1761.83", "547.71", "0.808"],
                ["Linear Regression", "1762.46", "547.99", "0.808"],
                ["Ridge Regression", "1762.94", "547.98", "0.808"],
                ["XGBoost", "2387.82", "534.66", "0.648"],
            ],
        },
        {
            "layout": "gallery",
            "kicker": "Model Diagnostics",
            "title": "Diagnostics and Business Utility (5.7, 5.8, 5.9.6)",
            "summary": [
                "Model ranking, residual diagnostics, and lift all support selection of a tuned regularized model.",
                "The top-decile lift demonstrates strong practical targeting value.",
            ],
            "images": [
                {"path": FIGURES / "model_comparison.png", "caption": "Model metric comparison"},
                {"path": FIGURES / "feature_importance.png", "caption": "Global feature importance"},
                {"path": FIGURES / "best_model_analysis.png", "caption": "Residual and fit diagnostics"},
                {"path": FIGURES / "lift_chart.png", "caption": "Lift chart (business impact)"},
            ],
        },
        {
            "layout": "bullets_with_image",
            "kicker": "Explainability",
            "title": "SHAP Explainability (5.9.6)",
            "bullets": [
                "SHAP confirms direction and magnitude of global drivers.",
                "Monetary and frequency behavior dominate prediction contributions.",
                "Interpretability supports stakeholder trust and campaign governance.",
                "Lift@10% = " + metrics["lift_10"] + "x over random targeting.",
            ],
            "image": FIGURES / "shap_summary.png",
        },
        {
            "layout": "section",
            "kicker": "Phase 2",
            "title": "CausalML Phase",
            "subtitle": "From prediction to intervention impact estimation",
        },
        {
            "layout": "split",
            "kicker": "Causal Setup",
            "title": "Target, Treatment, Controls",
            "left_title": "Definitions",
            "left_bullets": [
                "Outcome: log1p(CLV).",
                "Treatment: campaign vs control assignment.",
                "Controls: RFM + behavioral features.",
            ],
            "right_title": "Required Methods",
            "right_bullets": [
                "LRSRegressor (S-learner style).",
                "XGBTRegressor (T-learner style).",
                "Feature importance and CATE distribution analyzed in notebook 06.",
            ],
        },
        {
            "layout": "table",
            "kicker": "Causal Results",
            "title": "Average Treatment Effects",
            "headers": ["Model", "ATE (log1p CLV)", "CI Lower", "CI Upper"],
            "rows": [
                ["LRSRegressor", metrics["ate_lrs"], "-0.3783", "0.0213"],
                ["XGBTRegressor", metrics["ate_xgbt"], "-0.2575", "-0.0770"],
            ],
        },
        {
            "layout": "gallery",
            "kicker": "Causal Explainability",
            "title": "Uplift Distribution and Causal Importance",
            "summary": [
                "Treatment effects vary substantially by customer profile (heterogeneous CATE).",
                "Causal importance supports policy targeting beyond raw CLV ranking.",
            ],
            "images": [
                {"path": REPORTS / "causal_cate_distribution.png", "caption": "CATE distribution"},
                {"path": REPORTS / "causal_feature_importance.png", "caption": "Causal feature importance"},
            ],
        },
        {
            "layout": "bullets",
            "kicker": "Validity",
            "title": "Threats to Validity (5.9.7)",
            "bullets": [
                "Treatment assignment is observational/simulated, not randomized.",
                "Potential unobserved confounding and temporal effects.",
                "Outliers and historical policy effects can bias relationships.",
                "External validity may differ across geographies and verticals.",
            ],
        },
        {
            "layout": "bullets_with_image",
            "kicker": "MLOps",
            "title": "Launching, Monitoring, Maintenance (5.10)",
            "bullets": [
                "Production checks: schema validation, feature quality checks, artifact versioning.",
                "Monitoring: RMSE drift, CLV capture@K, feature drift via PSI.",
                "Alerting: warning at PSI >= 0.10, critical at PSI >= 0.25.",
                "Observed max PSI in validation run: " + metrics["max_psi"] + " (stable).",
            ],
            "image": FIGURES / "segment_dashboard.png",
        },
        {
            "layout": "bullets",
            "kicker": "Conclusion",
            "title": "Conclusions, Lessons, and Next Steps (5.9.8, 5.9.9)",
            "bullets": [
                "An end-to-end enterprise DS pipeline was implemented and validated.",
                "Predictive + explainable + causal layers improve decision quality.",
                "Next step: validate uplift on real campaign logs / A/B tests.",
                "This repository is submission-ready with reproducible artifacts.",
            ],
        },
        {
            "layout": "hero",
            "kicker": "Thank You",
            "title": "Questions and Discussion",
            "subtitle": (
                "Repository: " + REPO_LINK + "\n"
                "Othmane Zizi (261255341) | Fares Joni (261254593) | "
                "Tanmay Giri (261272443)"
            ),
            "footer": "Nerds With Attitude - INSY 674",
        },
    ]


HTML_CSS = """
@import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@400;500;700&family=Source+Sans+3:wght@400;600;700&display=swap');

:root{
  --ink:#0c1821;
  --ink-2:#1b2a41;
  --sky:#00a8e8;
  --mint:#2ec4b6;
  --sun:#ffb703;
  --sand:#f8f5f1;
  --paper:#ffffff;
  --muted:#6b7280;
}
*{box-sizing:border-box}
html,body{margin:0;padding:0;background:#0b1220}
body{font-family:'Source Sans 3',system-ui,sans-serif}
@page { size: 13.333in 7.5in; margin: 0; }
.slide{
  width:13.333in;height:7.5in;position:relative;overflow:hidden;
  page-break-after:always;background:var(--sand);color:var(--ink);
}
.slide::before{
  content:"";position:absolute;right:-1.2in;top:-1.1in;width:3in;height:3in;border-radius:50%;
  background:radial-gradient(circle at 30% 30%, rgba(46,196,182,.32), rgba(46,196,182,.06));
}
.slide::after{
  content:"";position:absolute;left:-1.4in;bottom:-1.4in;width:3.2in;height:3.2in;border-radius:50%;
  background:radial-gradient(circle at 65% 60%, rgba(0,168,232,.24), rgba(0,168,232,.02));
}
.hero{background:linear-gradient(135deg,#0c1821,#1b2a41);color:#f8fbff}
.hero .kicker,.hero .subtitle,.hero .footer{color:#d6e7f5}
.section{background:linear-gradient(120deg,#102539,#1b2a41);color:#f8fbff}
.wrap{position:relative;z-index:2;padding:.62in .72in .54in .72in;height:100%;display:flex;flex-direction:column}
.kicker{font-family:'Space Grotesk',sans-serif;font-size:13px;letter-spacing:2px;text-transform:uppercase;color:#236f96;font-weight:700}
.title{font-family:'Space Grotesk',sans-serif;font-weight:700;font-size:44px;line-height:1.08;margin:10px 0 8px 0}
.subtitle{font-size:20px;color:#4b5563;line-height:1.35;margin:0 0 18px 0;max-width:11in;white-space:pre-line}
.footer{margin-top:auto;font-size:15px;color:#374151}
.cards2{display:grid;grid-template-columns:1fr 1fr;gap:18px;margin-top:6px}
.card{
  background:linear-gradient(165deg, rgba(255,255,255,.97), rgba(248,252,255,.95));
  border:1px solid rgba(0,0,0,.08);
  border-radius:16px;padding:18px 18px 14px 18px;box-shadow:0 10px 24px rgba(10,20,40,.09)
}
.card h3{margin:0 0 8px 0;font-family:'Space Grotesk',sans-serif;font-size:24px}
.list{margin:0;padding:0;list-style:none;display:flex;flex-direction:column;gap:8px}
.list li{font-size:21px;line-height:1.28;display:flex;gap:10px}
.dot{color:var(--sky);font-weight:700}
.stats{display:grid;grid-template-columns:repeat(4,1fr);gap:14px;margin:8px 0 12px}
.stat{
  background:linear-gradient(160deg,#ffffff,#f2f8fd);
  border-radius:14px;border:1px solid rgba(0,0,0,.05);padding:14px 14px 10px;
}
.stat .v{font-family:'Space Grotesk',sans-serif;font-size:34px;font-weight:700;color:#0c516b}
.stat .l{font-size:12px;text-transform:uppercase;letter-spacing:1.2px;color:#4b5563}
.split{display:grid;grid-template-columns:1.15fr .85fr;gap:16px;align-items:stretch}
.image{
  width:100%;height:4.45in;object-fit:contain;background:#fff;
  border-radius:14px;border:1px solid rgba(0,0,0,.08);box-shadow:0 12px 30px rgba(2,8,20,.12)
}
.summary{gap:6px;margin:4px 0 8px}
.summary li{font-size:16px;line-height:1.22}
.gallery-grid{display:grid;gap:10px;flex:1;align-content:start;min-height:0}
.gallery-grid.grid-2{grid-template-columns:repeat(2,1fr)}
.gallery-grid.grid-4{grid-template-columns:repeat(2,1fr)}
.gallery-grid.grid-5{grid-template-columns:repeat(3,1fr)}
.gallery-grid.grid-5 .panel:nth-child(4){grid-column:1 / 2}
.gallery-grid.grid-5 .panel:nth-child(5){grid-column:2 / 4}
.panel{
  margin:0;background:#fff;border:1px solid rgba(0,0,0,.08);border-radius:12px;
  padding:6px 6px 8px;display:flex;flex-direction:column;min-height:0;box-shadow:0 8px 22px rgba(2,8,20,.08)
}
.gallery-img{width:100%;object-fit:contain;background:#fff;border-radius:8px;border:1px solid rgba(0,0,0,.08)}
.grid-2 .gallery-img{height:3.05in}
.grid-4 .gallery-img{height:1.78in}
.grid-5 .gallery-img{height:1.32in}
.panel figcaption{margin-top:5px;font-size:12px;line-height:1.2;color:#334155;text-align:center}
.image-missing{
  width:100%;height:1.4in;border-radius:8px;border:1px dashed rgba(100,116,139,.8);
  color:#475569;font-size:12px;display:flex;align-items:center;justify-content:center;background:#f8fafc
}
table{width:100%;border-collapse:collapse;font-size:20px;background:#fff;border-radius:14px;overflow:hidden;box-shadow:0 8px 24px rgba(2,8,20,.08)}
th,td{padding:12px 14px;border-bottom:1px solid #edf2f7;text-align:left}
th{font-family:'Space Grotesk',sans-serif;background:#102b40;color:#f8fbff;font-size:16px;text-transform:uppercase;letter-spacing:1.2px}
tr:nth-child(even) td{background:#f9fcff}
tr:last-child td{border-bottom:none}
.animate{animation:rise .5s ease-out}
@keyframes rise{from{transform:translateY(8px);opacity:.2}to{transform:translateY(0);opacity:1}}
"""


def _list_items(items: List[str]) -> str:
    return "".join(
        f"<li><span class='dot'>&bull;</span><span>{html.escape(item)}</span></li>" for item in items
    )


def _image_specs(slide: Dict[str, object]) -> List[Dict[str, object]]:
    specs: List[Dict[str, object]] = []
    for item in list(slide.get("images", [])):
        if isinstance(item, dict):
            path = Path(str(item.get("path", "")))
            caption = str(item.get("caption", path.name))
        else:
            path = Path(str(item))
            caption = path.stem.replace("_", " ").title()
        specs.append({"path": path, "caption": caption})
    return specs


def render_slide_html(slide: Dict[str, object]) -> str:
    layout = str(slide.get("layout"))
    kicker = html.escape(str(slide.get("kicker", "")))
    title = html.escape(str(slide.get("title", "")))
    subtitle = html.escape(str(slide.get("subtitle", "")))

    if layout == "hero":
        return f"""
<section class="slide hero animate">
  <div class="wrap" style="justify-content:center;text-align:center;">
    <div class="kicker">{kicker}</div>
    <h1 class="title" style="font-size:64px;max-width:11.8in;margin:14px auto 12px;">{title}</h1>
    <p class="subtitle" style="margin:0 auto 24px;max-width:10.2in;color:#dcebf6">{subtitle}</p>
    <div class="footer" style="color:#dcebf6">{html.escape(str(slide.get("footer", "")))}</div>
  </div>
</section>
"""

    if layout == "section":
        return f"""
<section class="slide section animate">
  <div class="wrap" style="justify-content:center;align-items:center;text-align:center;">
    <div class="kicker">{kicker}</div>
    <h2 class="title" style="font-size:62px;max-width:10.8in">{title}</h2>
    <p class="subtitle" style="color:#dcebf6">{subtitle}</p>
  </div>
</section>
"""

    if layout == "split":
        left_title = html.escape(str(slide.get("left_title", "")))
        right_title = html.escape(str(slide.get("right_title", "")))
        left_items = _list_items(list(slide.get("left_bullets", [])))
        right_items = _list_items(list(slide.get("right_bullets", [])))
        return f"""
<section class="slide animate">
  <div class="wrap">
    <div class="kicker">{kicker}</div>
    <h2 class="title">{title}</h2>
    <div class="cards2" style="margin-top:14px;">
      <article class="card"><h3>{left_title}</h3><ul class="list">{left_items}</ul></article>
      <article class="card"><h3>{right_title}</h3><ul class="list">{right_items}</ul></article>
    </div>
  </div>
</section>
"""

    if layout == "stats":
        stats_html = ""
        for label, value in list(slide.get("stats", [])):
            stats_html += (
                "<div class='stat'>"
                f"<div class='v'>{html.escape(str(value))}</div>"
                f"<div class='l'>{html.escape(str(label))}</div>"
                "</div>"
            )
        return f"""
<section class="slide animate">
  <div class="wrap">
    <div class="kicker">{kicker}</div>
    <h2 class="title">{title}</h2>
    <div class="stats">{stats_html}</div>
    <ul class="list">{_list_items(list(slide.get("bullets", [])))}</ul>
  </div>
</section>
"""

    if layout == "table":
        headers = list(slide.get("headers", []))
        rows = list(slide.get("rows", []))
        thead = "".join(f"<th>{html.escape(str(h))}</th>" for h in headers)
        body = ""
        for row in rows:
            body += "<tr>" + "".join(f"<td>{html.escape(str(v))}</td>" for v in row) + "</tr>"
        return f"""
<section class="slide animate">
  <div class="wrap">
    <div class="kicker">{kicker}</div>
    <h2 class="title">{title}</h2>
    <table><thead><tr>{thead}</tr></thead><tbody>{body}</tbody></table>
  </div>
</section>
"""

    if layout == "gallery":
        summary = list(slide.get("summary", []))
        summary_html = ""
        if summary:
            summary_html = f"<ul class='list summary'>{_list_items([str(x) for x in summary])}</ul>"

        images = _image_specs(slide)
        count = len(images)
        grid_class = "grid-2" if count <= 2 else ("grid-5" if count == 5 else "grid-4")
        panels = ""
        for spec in images:
            path = Path(str(spec["path"]))
            caption = html.escape(str(spec["caption"]))
            uri = img_data_uri(path)
            if uri:
                body = f"<img class='gallery-img' src='{uri}' alt='{caption}' />"
            else:
                body = f"<div class='image-missing'>Missing: {html.escape(path.name)}</div>"
            panels += f"<figure class='panel'>{body}<figcaption>{caption}</figcaption></figure>"
        return f"""
<section class="slide animate">
  <div class="wrap">
    <div class="kicker">{kicker}</div>
    <h2 class="title">{title}</h2>
    {summary_html}
    <div class="gallery-grid {grid_class}">{panels}</div>
  </div>
</section>
"""

    if layout == "bullets_with_image":
        image_path = Path(str(slide.get("image", "")))
        uri = img_data_uri(image_path)
        return f"""
<section class="slide animate">
  <div class="wrap">
    <div class="kicker">{kicker}</div>
    <h2 class="title">{title}</h2>
    <div class="split">
      <ul class="list">{_list_items(list(slide.get("bullets", [])))}</ul>
      <img class="image" src="{uri}" />
    </div>
  </div>
</section>
"""

    # Default bullets layout.
    return f"""
<section class="slide animate">
  <div class="wrap">
    <div class="kicker">{kicker}</div>
    <h2 class="title">{title}</h2>
    <ul class="list">{_list_items(list(slide.get("bullets", [])))}</ul>
  </div>
</section>
"""


def render_html(slides: List[Dict[str, object]]) -> str:
    sections = "\n".join(render_slide_html(s) for s in slides)
    return (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        "<meta name='viewport' content='width=device-width,initial-scale=1'>"
        f"<style>{HTML_CSS}</style></head><body>{sections}</body></html>"
    )


async def html_to_pdf(html_path: Path, pdf_path: Path) -> None:
    from playwright.async_api import async_playwright

    file_url = html_path.resolve().as_uri()
    try:
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page(viewport={"width": 1920, "height": 1080})
            await page.goto(file_url, wait_until="networkidle")
            await page.pdf(
                path=str(pdf_path),
                print_background=True,
                width="13.333in",
                height="7.5in",
                margin={"top": "0in", "right": "0in", "bottom": "0in", "left": "0in"},
            )
            await browser.close()
    except Exception:
        subprocess.run([sys.executable, "-m", "playwright", "install", "chromium"], check=True)
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page(viewport={"width": 1920, "height": 1080})
            await page.goto(file_url, wait_until="networkidle")
            await page.pdf(
                path=str(pdf_path),
                print_background=True,
                width="13.333in",
                height="7.5in",
                margin={"top": "0in", "right": "0in", "bottom": "0in", "left": "0in"},
            )
            await browser.close()


def _add_bg(slide, dark: bool = False) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    if dark:
        bg.fill.fore_color.rgb = RGBColor(16, 37, 57)
    else:
        bg.fill.fore_color.rgb = RGBColor(248, 245, 241)
    bg.line.fill.background()

    bubble1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(-1.0), Inches(3.0), Inches(3.0))
    bubble1.fill.solid()
    bubble1.fill.fore_color.rgb = RGBColor(46, 196, 182)
    bubble1.fill.transparency = 0.72
    bubble1.line.fill.background()

    bubble2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(5.8), Inches(3.2), Inches(3.2))
    bubble2.fill.solid()
    bubble2.fill.fore_color.rgb = RGBColor(0, 168, 232)
    bubble2.fill.transparency = 0.78
    bubble2.line.fill.background()


def _write_text(shape, text: str, size: int = 20, bold: bool = False, color: RGBColor = RGBColor(12, 24, 33),
                align: PP_ALIGN = PP_ALIGN.LEFT, font: str = "Avenir Next") -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color


def _add_bullet_box(slide, x, y, w, h, bullets: List[str], dark: bool = False) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Avenir Next"
        p.font.size = Pt(22)
        p.font.bold = False
        p.font.color.rgb = RGBColor(223, 237, 248) if dark else RGBColor(18, 34, 47)


def _add_summary_box(slide, summary: List[str], x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(summary):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        p.space_after = Pt(3)
        p.font.name = "Avenir Next"
        p.font.size = Pt(14)
        p.font.bold = False
        p.font.color.rgb = RGBColor(25, 45, 60)


def _gallery_slots(count: int, y_top: float) -> List[Tuple[float, float, float, float]]:
    if count <= 2:
        return [
            (0.7, y_top, 5.9, 4.35),
            (6.75, y_top, 5.9, 4.35),
        ][:count]
    if count == 5:
        return [
            (0.7, y_top, 3.85, 2.1),
            (4.75, y_top, 3.85, 2.1),
            (8.8, y_top, 3.85, 2.1),
            (0.7, y_top + 2.25, 5.9, 2.1),
            (6.75, y_top + 2.25, 5.9, 2.1),
        ]
    return [
        (0.7, y_top, 5.9, 2.25),
        (6.75, y_top, 5.9, 2.25),
        (0.7, y_top + 2.4, 5.9, 2.25),
        (6.75, y_top + 2.4, 5.9, 2.25),
    ][:count]


def _add_gallery_panel(slide, path: Path, caption: str, x: float, y: float, w: float, h: float) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.color.rgb = RGBColor(225, 234, 242)
    panel.line.width = Pt(1)

    img_x = x + 0.08
    img_y = y + 0.08
    img_w = max(0.1, w - 0.16)
    img_h = max(0.1, h - 0.48)

    if path.exists():
        # Fit image to panel without cropping; preserve aspect ratio and center it.
        try:
            from PIL import Image  # type: ignore

            with Image.open(path) as im:
                px_w, px_h = im.size
            if px_w > 0 and px_h > 0:
                box_ratio = img_w / img_h
                img_ratio = px_w / px_h
                if img_ratio > box_ratio:
                    draw_w = img_w
                    draw_h = img_w / img_ratio
                    draw_x = img_x
                    draw_y = img_y + (img_h - draw_h) / 2
                else:
                    draw_h = img_h
                    draw_w = img_h * img_ratio
                    draw_x = img_x + (img_w - draw_w) / 2
                    draw_y = img_y
            else:
                draw_x, draw_y, draw_w, draw_h = img_x, img_y, img_w, img_h
        except Exception:
            draw_x, draw_y, draw_w, draw_h = img_x, img_y, img_w, img_h

        slide.shapes.add_picture(
            str(path),
            Inches(draw_x),
            Inches(draw_y),
            width=Inches(draw_w),
            height=Inches(draw_h),
        )
    else:
        missing = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.2), Inches(w - 0.2), Inches(h - 0.7))
        _write_text(
            missing,
            f"Missing image: {path.name}",
            size=12,
            bold=False,
            color=RGBColor(71, 85, 105),
            align=PP_ALIGN.CENTER,
            font="Avenir Next",
        )

    cap = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + h - 0.33), Inches(w - 0.2), Inches(0.25))
    _write_text(
        cap,
        caption,
        size=11,
        bold=False,
        color=RGBColor(41, 59, 75),
        align=PP_ALIGN.CENTER,
        font="Avenir Next",
    )


def build_editable_pptx(slides: List[Dict[str, object]], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for s in slides:
        layout = str(s.get("layout"))
        dark = layout in {"hero", "section"}
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide, dark=dark)

        # Kicker
        kicker_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(5.5), Inches(0.4))
        _write_text(
            kicker_box,
            str(s.get("kicker", "")),
            size=12,
            bold=True,
            color=RGBColor(214, 231, 245) if dark else RGBColor(35, 111, 150),
            font="Avenir Next",
        )

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.9), Inches(1.2))
        _write_text(
            title_box,
            str(s.get("title", "")),
            size=46 if layout in {"hero", "section"} else 36,
            bold=True,
            color=RGBColor(248, 251, 255) if dark else RGBColor(12, 24, 33),
            font="Avenir Next",
        )

        if layout in {"hero", "section"}:
            sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6))
            _write_text(
                sub_box,
                str(s.get("subtitle", "")),
                size=22,
                bold=False,
                color=RGBColor(220, 235, 246),
                align=PP_ALIGN.CENTER,
                font="Avenir Next",
            )
            if s.get("footer"):
                footer = slide.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5))
                _write_text(
                    footer,
                    str(s.get("footer", "")),
                    size=14,
                    color=RGBColor(220, 235, 246),
                    align=PP_ALIGN.CENTER,
                    font="Avenir Next",
                )
            continue

        if layout == "split":
            left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.0), Inches(6.1), Inches(4.8))
            left.fill.solid()
            left.fill.fore_color.rgb = RGBColor(255, 255, 255)
            left.line.color.rgb = RGBColor(233, 240, 246)
            left.line.width = Pt(1)

            right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(2.0), Inches(5.65), Inches(4.8))
            right.fill.solid()
            right.fill.fore_color.rgb = RGBColor(255, 255, 255)
            right.line.color.rgb = RGBColor(233, 240, 246)
            right.line.width = Pt(1)

            left_title = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.5), Inches(0.5))
            _write_text(left_title, str(s.get("left_title", "")), size=24, bold=True)
            right_title = slide.shapes.add_textbox(Inches(7.25), Inches(2.2), Inches(4.9), Inches(0.5))
            _write_text(right_title, str(s.get("right_title", "")), size=24, bold=True)

            _add_bullet_box(slide, Inches(1.0), Inches(2.8), Inches(5.5), Inches(3.8), list(s.get("left_bullets", [])))
            _add_bullet_box(slide, Inches(7.25), Inches(2.8), Inches(4.9), Inches(3.8), list(s.get("right_bullets", [])))
            continue

        if layout == "stats":
            stats = list(s.get("stats", []))
            for i, (label, value) in enumerate(stats):
                x = 0.7 + i * 3.12
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.9), Inches(1.45))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card.line.color.rgb = RGBColor(233, 240, 246)
                card.line.width = Pt(1)
                v = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.15), Inches(2.5), Inches(0.65))
                _write_text(v, str(value), size=30, bold=True, color=RGBColor(12, 81, 107))
                l = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.85), Inches(2.5), Inches(0.35))
                _write_text(l, str(label), size=11, bold=True, color=RGBColor(75, 85, 99))

            _add_bullet_box(slide, Inches(0.9), Inches(3.7), Inches(12.0), Inches(2.7), list(s.get("bullets", [])))
            continue

        if layout == "table":
            headers = list(s.get("headers", []))
            rows = list(s.get("rows", []))
            table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.8)).table
            for j, h in enumerate(headers):
                c = table.cell(0, j)
                c.text = str(h)
                p = c.text_frame.paragraphs[0]
                p.font.name = "Avenir Next"
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = RGBColor(248, 251, 255)
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(16, 43, 64)
            for i, row in enumerate(rows, start=1):
                for j, v in enumerate(row):
                    c = table.cell(i, j)
                    c.text = str(v)
                    p = c.text_frame.paragraphs[0]
                    p.font.name = "Avenir Next"
                    p.font.size = Pt(13)
                    p.font.color.rgb = RGBColor(12, 24, 33)
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(249, 252, 255) if i % 2 == 0 else RGBColor(255, 255, 255)
            continue

        if layout == "gallery":
            summary = [str(x) for x in list(s.get("summary", []))]
            y_top = 1.95
            if summary:
                _add_summary_box(slide, summary, x=0.8, y=1.8, w=12.0, h=0.55)
                y_top = 2.15

            images = _image_specs(s)
            slots = _gallery_slots(len(images), y_top=y_top)
            for spec, (x, y, w, h) in zip(images, slots):
                _add_gallery_panel(
                    slide,
                    path=Path(str(spec["path"])),
                    caption=str(spec["caption"]),
                    x=x,
                    y=y,
                    w=w,
                    h=h,
                )
            continue

        # bullets / bullets_with_image default branch.
        has_image = layout == "bullets_with_image"
        bullet_w = Inches(7.6 if has_image else 12.0)
        _add_bullet_box(slide, Inches(0.8), Inches(2.0), bullet_w, Inches(4.8), list(s.get("bullets", [])))

        if has_image:
            image_path = Path(str(s.get("image", "")))
            if image_path.exists():
                box_x, box_y, box_w, box_h = 8.6, 2.0, 4.2, 4.8
                draw_x, draw_y, draw_w, draw_h = box_x, box_y, box_w, box_h
                try:
                    from PIL import Image  # type: ignore

                    with Image.open(image_path) as im:
                        px_w, px_h = im.size
                    if px_w > 0 and px_h > 0:
                        box_ratio = box_w / box_h
                        img_ratio = px_w / px_h
                        if img_ratio > box_ratio:
                            draw_w = box_w
                            draw_h = box_w / img_ratio
                            draw_x = box_x
                            draw_y = box_y + (box_h - draw_h) / 2
                        else:
                            draw_h = box_h
                            draw_w = box_h * img_ratio
                            draw_x = box_x + (box_w - draw_w) / 2
                            draw_y = box_y
                except Exception:
                    pass

                slide.shapes.add_picture(
                    str(image_path),
                    Inches(draw_x),
                    Inches(draw_y),
                    width=Inches(draw_w),
                    height=Inches(draw_h),
                )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def write_html(slides: List[Dict[str, object]], out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(render_html(slides), encoding="utf-8")


async def main() -> None:
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    metrics = load_dynamic_metrics()
    slides = build_slides(metrics)

    print("[1/3] Building HTML...")
    write_html(slides, HTML_OUT)
    print(f"  - {HTML_OUT}")

    print("[2/3] Converting HTML to PDF...")
    await html_to_pdf(HTML_OUT, PDF_OUT)
    print(f"  - {PDF_OUT}")

    print("[3/3] Converting HTML content to editable PPTX...")
    build_editable_pptx(slides, PPTX_OUT)
    print(f"  - {PPTX_OUT}")

    print("Done.")


if __name__ == "__main__":
    asyncio.run(main())
